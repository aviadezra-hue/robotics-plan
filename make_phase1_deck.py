"""
Generates 'Phase1-ROS2-Internalize.pptx' — a self-study PowerPoint deck
to help Aviad internalize Phase 1 of the robotics plan:
ROS 2 architecture, nodes, topics, services, parameters, launch, bags,
custom messages, and the publisher/subscriber coding pattern.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ------------- Style ---------------
ACCENT       = RGBColor(0x41, 0x56, 0xC8)   # deep indigo
ACCENT_SOFT  = RGBColor(0xE6, 0xEA, 0xFA)
FG           = RGBColor(0x1F, 0x29, 0x37)
MUTED        = RGBColor(0x55, 0x66, 0x7A)
CODE_BG      = RGBColor(0x0F, 0x17, 0x2A)
CODE_FG      = RGBColor(0xE2, 0xE8, 0xF0)
GOOD         = RGBColor(0x16, 0xA3, 0x4A)
WARN         = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ------------- Helpers ---------------
def add_slide():
    slide = prs.slides.add_slide(BLANK)
    # accent strip along the top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.25))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    return slide

def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=FG, align=PP_ALIGN.LEFT,
             font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top  = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb

def add_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.7),
             title, size=32, bold=True, color=ACCENT)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.5),
                 subtitle, size=16, color=MUTED)

def add_bullets(slide, left, top, width, height, items,
                size=18, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            head, rest = item
            r = p.add_run()
            r.text = head + " "
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = ACCENT; r.font.name = "Calibri"
            r = p.add_run()
            r.text = rest
            r.font.size = Pt(size); r.font.color.rgb = FG
            r.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.size = Pt(size); r.font.color.rgb = FG
            r.font.name = "Calibri"
    return tb

def add_code_block(slide, left, top, width, height, code, size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.line.color.rgb = ACCENT
    box.line.width = Pt(0.75)
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top  = Inches(0.12); tf.margin_bottom = Inches(0.12)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else " "
        r.font.size = Pt(size)
        r.font.name = "Consolas"
        r.font.color.rgb = CODE_FG
    return box

def add_callout(slide, left, top, width, height, text,
                color=ACCENT_SOFT, border=ACCENT, size=15):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.line.color.rgb = border
    box.line.width = Pt(1.5)
    box.fill.solid(); box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top  = Inches(0.12); tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.color.rgb = FG
    r.font.name = "Calibri"
    return box

def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             "Phase 1 — ROS 2 Fundamentals · Internalization Deck",
             size=10, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{page_no} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ------------ Slides ----------------
slides_meta = []  # (title, subtitle, build_fn)

def s_cover():
    s = add_slide()
    # big title block
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.2),
             "Phase 1: ROS 2 Fundamentals", size=48, bold=True, color=ACCENT)
    add_text(s, Inches(0.7), Inches(3.1), Inches(12), Inches(0.7),
             "Internalize the mental model — nodes, topics, services, parameters,",
             size=22, color=FG)
    add_text(s, Inches(0.7), Inches(3.55), Inches(12), Inches(0.7),
             "launch, bags, custom messages, and the pub/sub coding pattern.",
             size=22, color=FG)
    add_callout(s, Inches(0.7), Inches(5.0), Inches(12), Inches(1.4),
                "How to use this deck: spend 1–2 min per slide. After Slide 18, "
                "close the deck and try to recreate the talker/listener and the "
                "odom_demo package from memory. The slides you struggle to recall "
                "are the topics to review.", size=16)
    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
             "Aviad · ROS 2 Jazzy · Ubuntu 24.04 (WSL2)",
             size=14, color=MUTED, align=PP_ALIGN.LEFT)
slides_meta.append(("cover", s_cover))

def s_why_ros2():
    s = add_slide()
    add_title(s, "Why ROS 2 exists (in one slide)",
              "A middleware so robotics code is composable, language-agnostic, and replaceable.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5),
                [
                    ("Problem:", "robots = many concurrent processes (sensors, planners, controllers, UI). Each one needs to talk to others — over messages, with timing, across languages, across machines."),
                    ("ROS 2's answer:", "a pub/sub middleware (DDS) + a graph of small processes (nodes) + a contract language (.msg / .srv / .action) + a build system (colcon / ament)."),
                    ("Consequence:", "your wall-follower works whether the robot is a sim turtle, a TurtleBot, or a custom build — as long as the interface (/cmd_vel + /scan) matches."),
                    ("What you actually learn in Phase 1:", "the building blocks every Phase 2/3/4/5 demo is made of."),
                ], size=18)
slides_meta.append(("why", s_why_ros2))

def s_graph():
    s = add_slide()
    add_title(s, "The ROS 2 graph", "Everything is a node. Nodes talk over named channels.")
    # ASCII-style diagram via shapes
    # Three node ovals + arrows
    def node(x, y, label):
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(2.1), Inches(0.9))
        shp.fill.solid(); shp.fill.fore_color.rgb = ACCENT_SOFT
        shp.line.color.rgb = ACCENT; shp.line.width = Pt(2)
        tf = shp.text_frame; tf.margin_left = Inches(0.05); tf.margin_right=Inches(0.05)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT
        return shp

    def arrow(x1, y1, x2, y2, label):
        ln = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        ln.line.color.rgb = MUTED; ln.line.width = Pt(2)
        # arrow end
        ln.line.fill.solid()
        # label
        midx = (x1+x2)/2 - 0.6
        midy = (y1+y2)/2 - 0.25
        add_text(s, Inches(midx), Inches(midy), Inches(1.6), Inches(0.4),
                 label, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    node(1.0, 2.5, "talker\n(publisher)")
    node(5.6, 2.5, "/chatter\n(topic)")
    node(10.2, 2.5, "listener\n(subscriber)")
    arrow(3.1, 2.95, 5.6, 2.95, "publishes →")
    arrow(7.7, 2.95, 10.2, 2.95, "delivers →")

    add_bullets(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                [
                    ("Node:", "a process. Owns its state. Created via rclpy.Node (Python) or rclcpp::Node (C++)."),
                    ("Topic:", "a named, typed channel. Multiple publishers/subscribers allowed. Best-effort or reliable QoS."),
                    ("Discovery is automatic:", "every node on the same DOMAIN_ID sees every other node — no central broker."),
                ], size=17)
slides_meta.append(("graph", s_graph))

def s_topics_vs_services_vs_actions():
    s = add_slide()
    add_title(s, "Three communication patterns",
              "Pick the right one or your robot's behaviour will fight you.")
    # Three columns
    cols = [
        ("Topics", "🔁",
         "Continuous data streams.\n\nPub/sub, fire-and-forget.\nMany→many.\n\nExamples:\n• /scan (lidar @ 5 Hz)\n• /odom (state @ 50 Hz)\n• /cmd_vel (velocity setpoint)"),
        ("Services", "📞",
         "One-shot request/response.\n\nBlocking, synchronous.\nOne client → one server.\n\nExamples:\n• /spawn (create turtle)\n• /set_pen (change colour)\n• /add_two_ints"),
        ("Actions", "⏳",
         "Long-running goals with\nfeedback + cancel.\n\nGoal/result/feedback.\nOne client → one server.\n\nExamples:\n• /navigate_to_pose\n• /follow_path\n• /pick_object"),
    ]
    for i, (name, icon, body) in enumerate(cols):
        x = 0.5 + i * 4.3
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(1.7), Inches(4), Inches(5.2))
        card.fill.solid(); card.fill.fore_color.rgb = ACCENT_SOFT
        card.line.color.rgb = ACCENT; card.line.width = Pt(1.5)
        add_text(s, Inches(x+0.15), Inches(1.85), Inches(3.7), Inches(0.5),
                 f"{icon}  {name}", size=22, bold=True, color=ACCENT)
        add_text(s, Inches(x+0.15), Inches(2.5), Inches(3.7), Inches(4.2),
                 body, size=15, color=FG)
slides_meta.append(("3patterns", s_topics_vs_services_vs_actions))

def s_topics_deep():
    s = add_slide()
    add_title(s, "Topics in depth", "The bread and butter of ROS 2.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(4.5),
                [
                    ("Strongly typed:", "every topic has a message type. Mismatched types → silent no-delivery."),
                    ("QoS settings:", "reliability (best-effort vs reliable), durability (volatile vs transient_local), depth, deadline. Defaults work for most things."),
                    ("Multiple publishers OK:", "first writer wins by timestamp. Subscribers don't care who wrote."),
                    ("Topic remapping at launch:", "rename a node's /cmd_vel → /turtle1/cmd_vel without touching the node's code."),
                ], size=16)
    add_code_block(s, Inches(8.3), Inches(1.7), Inches(4.5), Inches(4.5),
                   "ros2 topic list\n"
                   "ros2 topic info /odom\n"
                   "ros2 topic echo /odom\n"
                   "ros2 topic hz /odom\n"
                   "ros2 topic pub /cmd_vel \\\n"
                   "  geometry_msgs/msg/Twist \\\n"
                   "  '{linear: {x: 0.2},\n"
                   "    angular: {z: 0.5}}'\n"
                   "ros2 topic pub /cmd_vel ... -r 1\n"
                   "  # republish at 1 Hz",
                   size=13)
slides_meta.append(("topics", s_topics_deep))

def s_services_params():
    s = add_slide()
    add_title(s, "Services & parameters",
              "Synchronous calls + live-tunable config.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(5),
                [
                    ("Services = call/reply", "via .srv files. Examples: spawn turtle, kill turtle, set_pen."),
                    ("Use them sparingly:", "topics scale better. Services block the caller until the server replies."),
                    ("Parameters = live config", "per-node typed values. Declared in code, set at launch, can be changed at runtime."),
                    ("Use cases:", "publish rate, frame_id, max velocity, K-gains, file paths."),
                    ("Why this matters:", "tuning Nav2 in Phase 3 is 80% twiddling parameters — you'll do this every day."),
                ], size=16)
    add_code_block(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5),
                   "# Services\n"
                   "ros2 service list\n"
                   "ros2 service type /spawn\n"
                   "ros2 service call /spawn \\\n"
                   "  turtlesim/srv/Spawn \\\n"
                   "  '{x: 5, y: 5, name: t2}'\n"
                   "\n"
                   "# Parameters\n"
                   "ros2 param list /turtlesim\n"
                   "ros2 param get /turtlesim background_r\n"
                   "ros2 param set /turtlesim background_r 200\n"
                   "ros2 param dump /turtlesim > tsim.yaml",
                   size=13)
slides_meta.append(("services", s_services_params))

def s_launch():
    s = add_slide()
    add_title(s, "Launch files = the deployment unit",
              "One .launch.py starts N nodes with the right params and remappings.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(6.3), Inches(5),
                [
                    ("Why:", "real systems aren't one node — Nav2 = ~15 nodes, MoveIt = ~10."),
                    ("Written in Python", "(launch_ros + launch). Composable: launch files can include other launch files."),
                    ("DeclareLaunchArgument", "exposes CLI args: ros2 launch pkg foo.launch.py rate:=50.0"),
                    ("Remappings:", "rename topics per-node without touching code."),
                    ("Parameters:", "passed in as a dict — overrides node defaults."),
                ], size=16)
    add_code_block(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5),
                   "from launch import LaunchDescription\n"
                   "from launch_ros.actions import Node\n"
                   "\n"
                   "def generate_launch_description():\n"
                   "  return LaunchDescription([\n"
                   "    Node(\n"
                   "      package='odom_demo',\n"
                   "      executable='odom_publisher',\n"
                   "      name='odom_pub',\n"
                   "      parameters=[{'publish_rate_hz': 10.0}],\n"
                   "      remappings=[('/odom','/turtle1/odom')],\n"
                   "    ),\n"
                   "  ])",
                   size=12)
slides_meta.append(("launch", s_launch))

def s_bags():
    s = add_slide()
    add_title(s, "Bags = record / replay reality",
              "rosbag2 is your debugger, your dataset, and your CI input.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(5),
                [
                    ("Record any/all topics:", "ros2 bag record /odom /scan"),
                    ("Replay later:", "ros2 bag play <dir> — every subscriber sees the original timing."),
                    ("Use cases:", "debug a crash from yesterday's run, build a regression test, train a perception model offline, share a reproducer."),
                    ("Stored as MCAP/SQLite", "with full QoS + type info — fully self-contained, no schema setup needed on replay."),
                    ("Pro tip:", "log bags from every real-robot run. Storage is cheap; missing data is not."),
                ], size=17)
    add_code_block(s, Inches(8.3), Inches(1.7), Inches(4.5), Inches(2.6),
                   "ros2 bag record -a\n"
                   "ros2 bag record /odom /scan\n"
                   "ros2 bag info <dir>\n"
                   "ros2 bag play <dir> --rate 2.0",
                   size=13)
slides_meta.append(("bags", s_bags))

def s_custom_msg():
    s = add_slide()
    add_title(s, "Custom messages (.msg)",
              "Define your own typed contract once → use from Python and C++.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(6.5), Inches(5),
                [
                    ("Lives in a separate package", "of build type ament_cmake (Python build can't generate IDL)."),
                    ("Use rosidl_generate_interfaces", "in CMakeLists.txt; declare rosidl_default_generators + runtime in package.xml."),
                    ("Field types:", "primitive (float64, int32, string), arrays (float64[5]), other msgs (geometry_msgs/Pose), builtin_interfaces/Time."),
                    ("Why:", "your team can share a precise data contract — no JSON guessing, no protobuf setup. ROS 2 builds the Python + C++ classes for you."),
                ], size=16)
    add_code_block(s, Inches(7.2), Inches(1.7), Inches(5.6), Inches(5),
                   "# msg/VelocityStats.msg\n"
                   "builtin_interfaces/Time stamp\n"
                   "float64 linear_speed       # m/s\n"
                   "float64 angular_speed      # rad/s\n"
                   "float64 distance_traveled  # m\n"
                   "uint32  sample_count\n"
                   "\n"
                   "# Use from Python:\n"
                   "from odom_demo_msgs.msg \\\n"
                   "  import VelocityStats",
                   size=13)
slides_meta.append(("msg", s_custom_msg))

def s_pubsub_pattern():
    s = add_slide()
    add_title(s, "The pub/sub coding pattern",
              "Memorize this shape. Every Python node you'll ever write looks like it.")
    add_code_block(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(5.0),
                   "import rclpy\n"
                   "from rclpy.node import Node\n"
                   "from std_msgs.msg import String\n"
                   "\n"
                   "class MyNode(Node):\n"
                   "  def __init__(self):\n"
                   "    super().__init__('my_node')\n"
                   "    self.declare_parameter('rate_hz', 10.0)\n"
                   "    rate = self.get_parameter('rate_hz').value\n"
                   "    self.pub = self.create_publisher(String, '/out', 10)\n"
                   "    self.sub = self.create_subscription(\n"
                   "        String, '/in', self.on_msg, 10)\n"
                   "    self.timer = self.create_timer(1.0/rate, self.tick)\n"
                   "\n"
                   "  def on_msg(self, msg): ...\n"
                   "  def tick(self):         ...\n"
                   "\n"
                   "def main():\n"
                   "  rclpy.init(); rclpy.spin(MyNode()); rclpy.shutdown()",
                   size=13)
    add_bullets(s, Inches(8.3), Inches(1.7), Inches(4.5), Inches(5),
                [
                    ("Subclass Node", "→ get logging, params, pubs, subs, timers, services."),
                    ("Declare parameters first", "with defaults."),
                    ("Create pubs / subs / timers in __init__", ", do work in callbacks."),
                    ("Never block in a callback", "— it stops other callbacks."),
                    ("Use rclpy.spin", "to run the event loop."),
                ], size=14)
slides_meta.append(("pattern", s_pubsub_pattern))

def s_workspace():
    s = add_slide()
    add_title(s, "Workspace & package anatomy",
              "Where code lives, how it's built, why colcon scans src/ recursively.")
    add_code_block(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(5.2),
                   "~/ros2_ws/                ← workspace root\n"
                   "├── src/                  ← your source goes here\n"
                   "│   ├── odom_demo_msgs/   ← ament_cmake (for .msg)\n"
                   "│   │   ├── CMakeLists.txt\n"
                   "│   │   ├── package.xml\n"
                   "│   │   └── msg/VelocityStats.msg\n"
                   "│   └── odom_demo/        ← ament_python (for nodes)\n"
                   "│       ├── package.xml\n"
                   "│       ├── setup.py\n"
                   "│       ├── odom_demo/    ← INNER python module\n"
                   "│       │   ├── __init__.py\n"
                   "│       │   ├── odom_publisher.py\n"
                   "│       │   └── odom_subscriber.py\n"
                   "│       └── launch/demo.launch.py\n"
                   "├── build/    ← colcon scratch (gitignored)\n"
                   "├── install/  ← built artefacts (source setup.bash)\n"
                   "└── log/      ← per-build logs",
                   size=12)
    add_bullets(s, Inches(8.3), Inches(1.7), Inches(4.5), Inches(5),
                [
                    ("colcon build", "from the workspace root scans src/ recursively and builds in dependency order."),
                    ("source install/setup.bash", "puts the built packages on your PATH + PYTHONPATH for this shell."),
                    ("--packages-select", "builds only what you name (much faster)."),
                    ("Re-build after editing .msg", "→ always. After editing .py → optional with symlink-install."),
                ], size=13)
slides_meta.append(("workspace", s_workspace))

def s_cli_cheatsheet():
    s = add_slide()
    add_title(s, "Phase 1 CLI cheat sheet",
              "If you forget any of these, run the command with --help.")
    # Two-column code blocks
    add_code_block(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.5),
                   "# Discovery / introspection\n"
                   "ros2 node    list / info /name\n"
                   "ros2 topic   list / info / type\n"
                   "ros2 service list / type / call\n"
                   "ros2 param   list / get / set / dump\n"
                   "ros2 action  list / info\n"
                   "ros2 interface show <pkg>/msg/<Name>\n"
                   "\n"
                   "# Live traffic\n"
                   "ros2 topic echo /odom\n"
                   "ros2 topic hz /odom\n"
                   "ros2 topic pub /cmd_vel ... -r 1\n"
                   "\n"
                   "# Visualize the graph\n"
                   "rqt_graph",
                   size=13)
    add_code_block(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(5.5),
                   "# Build / source\n"
                   "cd ~/ros2_ws\n"
                   "colcon build --packages-select <pkg>\n"
                   "source install/setup.bash\n"
                   "\n"
                   "# Scaffold packages\n"
                   "ros2 pkg create --build-type ament_python \\\n"
                   "  <name> --license Apache-2.0 \\\n"
                   "  --dependencies rclpy ...\n"
                   "\n"
                   "# Launch & bags\n"
                   "ros2 launch <pkg> <file>.launch.py\n"
                   "ros2 bag record -a\n"
                   "ros2 bag play <dir>\n"
                   "\n"
                   "# Daemon (when discovery acts weird)\n"
                   "ros2 daemon stop && ros2 daemon start",
                   size=13)
slides_meta.append(("cheat", s_cli_cheatsheet))

def s_dataflow():
    s = add_slide()
    add_title(s, "Step 2.8 in one picture",
              "Mimic + 2 remappings — the whole secret of ROS 2 composability.")
    def node(x, y, w, label, color=ACCENT_SOFT):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(0.85))
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.color.rgb = ACCENT; shp.line.width = Pt(1.5)
        tf = shp.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT
        return shp
    def label(x, y, w, txt):
        add_text(s, Inches(x), Inches(y), Inches(w), Inches(0.4),
                 txt, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    def arr(x1, y, x2):
        ln = s.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
        ln.line.color.rgb = MUTED; ln.line.width = Pt(2)

    y_top = 2.1
    node(0.5, y_top, 2.2, "teleop_key")
    label(2.7, y_top+0.05, 1.4, "/turtle1/cmd_vel →")
    node(4.2, y_top, 3.2, "turtlesim_node\n(turtle1 + turtle2)")
    label(7.4, y_top+0.05, 1.6, "/turtle1/pose →")
    node(9.1, y_top, 2.2, "mimic")
    label(7.4, y_top+0.6, 1.6, "← /turtle2/cmd_vel")

    arr(2.7, y_top+0.42, 4.2)
    arr(7.4, y_top+0.42, 9.1)
    arr(9.1, y_top+1.0, 7.4)  # mimic's output back to sim
    add_bullets(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3),
                [
                    ("Mimic is a generic node.", "It listens on /input/pose, computes a twist, publishes on /output/cmd_vel."),
                    ("Remappings in the launch file", "rewire /input → /turtle1/pose and /output → /turtle2/cmd_vel."),
                    ("Result:", "turtle2 follows turtle1 with zero code changes to mimic."),
                    ("This is exactly", "how Nav2 wires 15 nodes together. Same pattern, bigger graph."),
                ], size=15)
slides_meta.append(("dataflow", s_dataflow))

def s_recall():
    s = add_slide()
    add_title(s, "Recall test — try this without peeking",
              "If you struggle on any of these, re-read that step.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5),
                [
                    ("1.", "What's the difference between a topic, a service, and an action?"),
                    ("2.", "Write the 6-line skeleton of a Python ROS 2 node (no imports needed)."),
                    ("3.", "Which package type (ament_cmake / ament_python) hosts .msg files? Why?"),
                    ("4.", "How do you change a running node's parameter without restarting it?"),
                    ("5.", "What does `colcon build --packages-select foo` do that bare `colcon build` doesn't?"),
                    ("6.", "Name the file you create in a launch file to expose `rate:=50.0` as a CLI override."),
                    ("7.", "What does `source install/setup.bash` actually modify in your shell?"),
                    ("8.", "Why does the same teleop work for turtlesim AND TurtleBot 3 AND real hardware?"),
                ], size=17)
slides_meta.append(("recall", s_recall))

def s_pitfalls():
    s = add_slide()
    add_title(s, "Phase 1 pitfalls — pre-empt these",
              "Every one of these has cost me at least 20 minutes the first time.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5),
                [
                    ("Forgot to source setup.bash", "in a new terminal → ros2 commands work but your custom pkg is invisible. Fix: it's in ~/.bashrc already (Step 1.6)."),
                    ("ros2 pkg create with absolute path", "skips the inner Python module folder. Always cd ~/ros2_ws/src first, then use the bare package name."),
                    ("Edited a .msg but didn't colcon build", "→ subscriber gets stale typedef → silent no-delivery."),
                    ("Two nodes with the same name on same DOMAIN_ID", "→ undefined behaviour. Use `__node:=` remapping if you need two instances."),
                    ("Topic exists but echo shows nothing", "→ QoS mismatch. Check `ros2 topic info -v /topic`."),
                    ("Discovery flaky after lots of starts/stops", "→ ros2 daemon stop && ros2 daemon start."),
                ], size=15)
slides_meta.append(("pitfalls", s_pitfalls))

def s_phase1_arc():
    s = add_slide()
    add_title(s, "Phase 1 in three arcs",
              "What you built, in order — and what each step unlocked.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5),
                [
                    ("Step 1 — Install + smoke test:", "got ROS 2 Jazzy on Ubuntu 24.04 (WSL2), proved pub/sub works end-to-end with talker/listener. → You now have a working dev environment."),
                    ("Step 2 — Explore the CLI with turtlesim:", "8 substeps drilling nodes, topics, services, parameters, launch files, bags, multi-node configs, and node-composition via mimic. → You now know how to INTROSPECT any ROS 2 system."),
                    ("Step 3 — Mini-project (odom_demo):", "scaffolded two packages (msgs + python), wrote a publisher + subscriber, defined a custom message, wired entry points, built with colcon, validated with a launch file + CLI overrides. → You now know how to WRITE any ROS 2 system."),
                    ("Whole loop:", "build → launch → introspect → tweak params → re-launch. This is the dev cycle for the rest of your robotics career."),
                ], size=16)
slides_meta.append(("arc", s_phase1_arc))

def s_mental_models():
    s = add_slide()
    add_title(s, "Five mental models worth keeping",
              "If you remember only these, you've internalized Phase 1.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5),
                [
                    ("1. Robots are graphs of small processes.", "Every behaviour you'll ever write is N nodes connected by typed channels."),
                    ("2. Interfaces are contracts.", ".msg / .srv / .action files ARE the API. Stable interface → swappable implementation."),
                    ("3. Launch files are the deployment.", "If you can launch it, you can ship it. Hand-running ros2 run is for debugging."),
                    ("4. CLI introspection beats reading code.", "ros2 topic / node / param / interface tells you what's ACTUALLY happening at runtime."),
                    ("5. Same workflow at every scale.", "turtlesim, TurtleBot 3, Nav2, your own robot — pub/sub + launch + params, just more of them."),
                ], size=17)
slides_meta.append(("models", s_mental_models))

def s_next():
    s = add_slide()
    add_title(s, "Next: Phase 2 — Simulation Deep Dive",
              "You graduate from turtle 2D to a physically simulated robot.")
    add_bullets(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.5),
                [
                    ("TurtleBot 3 Burger in Gazebo Classic", "— real lidar, real odometry, real friction."),
                    ("Same skills you just learned,", "applied to a robot that mirrors what real hardware exposes (/cmd_vel, /odom, /scan)."),
                    ("Mini-project:", "front_distance_monitor — a Python node that watches lidar and warns near walls. Same pub/sub pattern from Phase 1 Step 3, smarter math."),
                    ("Then Phase 3:", "Nav2 — give the robot a map and a goal, watch it drive itself there."),
                ], size=17)
    add_callout(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9),
                "Reminder: the resume marker on the site lives at "
                "Phase 1 Step 3.8. Click the 📍 floating pill on the live site "
                "to jump back to where you left off.",
                size=15)
slides_meta.append(("next", s_next))

# Build & save
total = len(slides_meta)
for i, (_, fn) in enumerate(slides_meta, 1):
    fn()

# Add footers (skip cover)
for i, slide in enumerate(prs.slides, 1):
    if i == 1:
        continue
    add_footer(slide, i, total)

out = r"C:\Users\avezra\repro-robotics-plan\Phase1-ROS2-Internalize.pptx"
prs.save(out)
print(f"Wrote {out} with {total} slides.")
